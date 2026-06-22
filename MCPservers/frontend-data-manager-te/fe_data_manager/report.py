"""PowerPoint report generation for a basictype dataset.

Charts are rendered with matplotlib (Agg backend) to temporary PNG files and embedded
into a ``.pptx`` built with python-pptx. If a reference/template ``.pptx`` is available it
is used for theme/slide-size, with its sample slides removed.
"""
from __future__ import annotations

import logging
import tempfile
from datetime import date, datetime
from pathlib import Path

import matplotlib

matplotlib.use("Agg")
import matplotlib.pyplot as plt  # noqa: E402
import numpy as np  # noqa: E402
from pptx import Presentation  # noqa: E402
from pptx.util import Emu, Pt  # noqa: E402
from pptx.dml.color import RGBColor  # noqa: E402
from pptx.enum.text import PP_ALIGN  # noqa: E402

from . import analysis, config  # noqa: E402

log = logging.getLogger(__name__)

_LIMIT_COLOR = "#c00000"
_HEADER_FILL = RGBColor(0x1F, 0x49, 0x7D)


# --------------------------------------------------------------------------- charts
def _fmt(value: float | None, digits: int = 4) -> str:
    if value is None or (isinstance(value, float) and np.isnan(value)):
        return "n/a"
    if value != 0 and (abs(value) < 1e-3 or abs(value) >= 1e5):
        return f"{value:.3e}"
    return f"{value:.{digits}g}"


def _boxplot_png(summary: "analysis.TestSummary", out: Path) -> Path | None:
    data = [w.values for w in summary.per_wafer if w.values.size]
    labels = [f"#{w.wafer_number}" for w in summary.per_wafer if w.values.size]
    if not data:
        return None
    fig, ax = plt.subplots(figsize=(5.2, 3.9), dpi=150)
    ax.boxplot(data, tick_labels=labels, showfliers=True, widths=0.6)
    if summary.usl is not None:
        ax.axhline(summary.usl, color=_LIMIT_COLOR, ls="--", lw=1.2, label="USL")
    if summary.lsl is not None:
        ax.axhline(summary.lsl, color=_LIMIT_COLOR, ls="-.", lw=1.2, label="LSL")
    ax.set_title(f"Boxplot by wafer — {summary.name}", fontsize=10)
    ax.set_xlabel("Wafer")
    ax.set_ylabel(summary.unit or "value")
    if summary.usl is not None or summary.lsl is not None:
        ax.legend(fontsize=8, loc="best")
    ax.grid(axis="y", ls=":", alpha=0.4)
    fig.tight_layout()
    fig.savefig(out)
    plt.close(fig)
    return out


def _cumfreq_png(summary: "analysis.TestSummary", out: Path) -> Path | None:
    plotted = False
    fig, ax = plt.subplots(figsize=(5.2, 3.9), dpi=150)
    cmap = plt.get_cmap("tab10")
    for i, w in enumerate(summary.per_wafer):
        if not w.values.size:
            continue
        xs = np.sort(w.values)
        ys = np.arange(1, xs.size + 1) / xs.size * 100.0
        ax.plot(xs, ys, ls=":", marker="", lw=1.4, color=cmap(i % 10),
                label=f"#{w.wafer_number}")
        plotted = True
    if not plotted:
        plt.close(fig)
        return None
    if summary.usl is not None:
        ax.axvline(summary.usl, color=_LIMIT_COLOR, ls="--", lw=1.2, label="USL")
    if summary.lsl is not None:
        ax.axvline(summary.lsl, color=_LIMIT_COLOR, ls="-.", lw=1.2, label="LSL")
    ax.set_title(f"Cumulative frequency — {summary.name}", fontsize=10)
    ax.set_xlabel(summary.unit or "value")
    ax.set_ylabel("Cumulative %")
    ax.set_ylim(0, 100)
    ax.legend(fontsize=7, ncol=2, loc="best")
    ax.grid(ls=":", alpha=0.4)
    fig.tight_layout()
    fig.savefig(out)
    plt.close(fig)
    return out


def _yield_png(yields: list["analysis.WaferYield"], out: Path) -> Path:
    labels = [f"#{y.wafer_number}" for y in yields]
    pct = [y.yield_pct for y in yields]
    fig, ax = plt.subplots(figsize=(7.5, 4.0), dpi=150)
    bars = ax.bar(labels, pct, color="#4472c4")
    for b, p in zip(bars, pct):
        ax.text(b.get_x() + b.get_width() / 2, p, f"{p:.1f}%",
                ha="center", va="bottom", fontsize=8)
    ax.set_title("Overall yield by wafer", fontsize=11)
    ax.set_xlabel("Wafer")
    ax.set_ylabel("Yield %")
    ax.set_ylim(0, 105)
    ax.grid(axis="y", ls=":", alpha=0.4)
    fig.tight_layout()
    fig.savefig(out)
    plt.close(fig)
    return out


def _pareto_png(pareto: list[tuple[int, int]], out: Path) -> Path | None:
    if not pareto:
        return None
    bins = [f"HBIN {b}" for b, _ in pareto]
    counts = [c for _, c in pareto]
    total = sum(counts) or 1
    cum = np.cumsum(counts) / total * 100.0
    fig, ax = plt.subplots(figsize=(7.5, 4.0), dpi=150)
    ax.bar(bins, counts, color="#ed7d31")
    ax.set_ylabel("Fail count")
    ax.set_xlabel("Hardware bin")
    ax.set_title("HBIN Pareto (fails)", fontsize=11)
    ax2 = ax.twinx()
    ax2.plot(bins, cum, color="#c00000", marker="o", lw=1.5)
    ax2.set_ylabel("Cumulative %")
    ax2.set_ylim(0, 105)
    ax.tick_params(axis="x", rotation=30)
    fig.tight_layout()
    fig.savefig(out)
    plt.close(fig)
    return out


def _trend_png(
    wafer_numbers: list[int], top_bins: list[int], trend: dict[int, list[float]], out: Path
) -> Path | None:
    if not top_bins:
        return None
    labels = [f"#{n}" for n in wafer_numbers]
    fig, ax = plt.subplots(figsize=(7.5, 4.0), dpi=150)
    cmap = plt.get_cmap("tab10")
    for i, b in enumerate(top_bins):
        ax.plot(labels, trend[b], marker="o", lw=1.5, color=cmap(i % 10),
                label=f"HBIN {b}")
    ax.set_title("HBIN trend across wafers", fontsize=11)
    ax.set_xlabel("Wafer")
    ax.set_ylabel("Fail %")
    ax.legend(fontsize=8)
    ax.grid(ls=":", alpha=0.4)
    fig.tight_layout()
    fig.savefig(out)
    plt.close(fig)
    return out


# ----------------------------------------------------------------------- pptx utils
def _new_presentation() -> Presentation:
    """Create a blank presentation, inheriting the slide size from the reference deck.

    We deliberately do not use the reference ``.pptx`` as the *base* file: removing its
    sample slides reliably leaves orphaned parts that corrupt the output. Instead we start
    blank and copy the reference slide dimensions so the geometry matches.
    """
    prs = Presentation()
    template = config.resolve_template()
    if template is not None:
        try:
            ref = Presentation(str(template))
            if ref.slide_width and ref.slide_height:
                prs.slide_width = ref.slide_width
                prs.slide_height = ref.slide_height
        except Exception as exc:  # noqa: BLE001
            log.warning("Could not read template %s: %s", template, exc)
    return prs


def _blank_layout(prs: Presentation):
    # Prefer a truly blank layout; fall back to the last available layout.
    layouts = prs.slide_layouts
    for layout in layouts:
        if layout.name and "blank" in layout.name.lower():
            return layout
    return layouts[6] if len(layouts) > 6 else layouts[-1]


def _add_textbox(slide, left, top, width, height, text, size=14, bold=False,
                 color=None, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    para = tf.paragraphs[0]
    para.alignment = align
    run = para.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    if color is not None:
        run.font.color.rgb = color
    return box


def _add_picture_fit(slide, png: Path, left, top, max_w, max_h):
    # python-pptx keeps the aspect ratio when only width is supplied.
    slide.shapes.add_picture(str(png), left, top, width=int(max_w))


# ------------------------------------------------------------------------- slides
def _slide(prs, layout):
    return prs.slides.add_slide(layout)


def _add_title_slide(prs, ds: "analysis.BasicTypeDataset", n_files: int):
    layout = _blank_layout(prs)
    slide = _slide(prs, layout)
    sw, sh = prs.slide_width, prs.slide_height
    _add_textbox(slide, Emu(int(sw * 0.08)), Emu(int(sh * 0.30)),
                 Emu(int(sw * 0.84)), Emu(int(sh * 0.20)),
                 f"{ds.p_code} \u2013 {ds.vf_code}\nFE Testing Results",
                 size=34, bold=True, align=PP_ALIGN.CENTER, color=_HEADER_FILL)
    subtitle = (
        f"Basictype: {ds.basictype}    |    Wafers: {len(ds.wafers)}    |    "
        f"EFF files: {n_files}\nTest numbers: {len(ds.test_numbers)}    |    "
        f"Generated: {date.today():%Y-%m-%d}"
    )
    _add_textbox(slide, Emu(int(sw * 0.08)), Emu(int(sh * 0.55)),
                 Emu(int(sw * 0.84)), Emu(int(sh * 0.15)),
                 subtitle, size=14, align=PP_ALIGN.CENTER)


def _add_chart_slide(prs, title: str, png: Path | None, note: str | None = None):
    layout = _blank_layout(prs)
    slide = _slide(prs, layout)
    sw, sh = prs.slide_width, prs.slide_height
    _add_textbox(slide, Emu(int(sw * 0.05)), Emu(int(sh * 0.03)),
                 Emu(int(sw * 0.90)), Emu(int(sh * 0.10)),
                 title, size=22, bold=True, color=_HEADER_FILL)
    if png is not None and png.exists():
        _add_picture_fit(slide, png, Emu(int(sw * 0.10)), Emu(int(sh * 0.16)),
                         int(sw * 0.80), int(sh * 0.78))
    if note:
        _add_textbox(slide, Emu(int(sw * 0.05)), Emu(int(sh * 0.90)),
                     Emu(int(sw * 0.90)), Emu(int(sh * 0.08)), note, size=11)
    return slide


def _add_yield_slide(prs, ds, yields, png):
    slide = _add_chart_slide(prs, "Overall Yield", png)
    sw, sh = prs.slide_width, prs.slide_height
    rows = len(yields) + 1
    table = slide.shapes.add_table(
        rows, 4,
        Emu(int(sw * 0.70)), Emu(int(sh * 0.16)),
        Emu(int(sw * 0.27)), Emu(int(sh * 0.06 * rows)),
    ).table
    for j, head in enumerate(("Wafer", "Total", "Pass", "Yield%")):
        table.cell(0, j).text = head
    for i, y in enumerate(yields, start=1):
        table.cell(i, 0).text = f"#{y.wafer_number}"
        table.cell(i, 1).text = str(y.total)
        table.cell(i, 2).text = str(y.passed)
        table.cell(i, 3).text = f"{y.yield_pct:.2f}"
    _shrink_table_font(table, 9)


def _add_test_slide(prs, summary: "analysis.TestSummary"):
    layout = _blank_layout(prs)
    slide = _slide(prs, layout)
    sw, sh = prs.slide_width, prs.slide_height

    title = f"Test {summary.number}: {summary.name}  [{summary.unit or '-'}]"
    _add_textbox(slide, Emu(int(sw * 0.04)), Emu(int(sh * 0.02)),
                 Emu(int(sw * 0.92)), Emu(int(sh * 0.07)),
                 title, size=20, bold=True, color=_HEADER_FILL)

    header = (
        f"USL: {_fmt(summary.usl)}    LSL: {_fmt(summary.lsl)}    "
        f"Overall Mean: {_fmt(summary.overall_mean)}"
    )
    _add_textbox(slide, Emu(int(sw * 0.04)), Emu(int(sh * 0.09)),
                 Emu(int(sw * 0.92)), Emu(int(sh * 0.05)), header, size=12, bold=True)

    return slide


def _add_test_charts(slide, prs, box_png: Path | None, cum_png: Path | None):
    sw, sh = prs.slide_width, prs.slide_height
    top = Emu(int(sh * 0.16))
    if box_png is not None and box_png.exists():
        _add_picture_fit(slide, box_png, Emu(int(sw * 0.02)), top,
                         int(sw * 0.47), int(sh * 0.62))
    if cum_png is not None and cum_png.exists():
        _add_picture_fit(slide, cum_png, Emu(int(sw * 0.50)), top,
                         int(sw * 0.47), int(sh * 0.62))


def _add_test_table(slide, prs, summary: "analysis.TestSummary"):
    sw, sh = prs.slide_width, prs.slide_height
    wafers = summary.per_wafer
    cols = len(wafers) + 1
    table = slide.shapes.add_table(
        4, cols,
        Emu(int(sw * 0.04)), Emu(int(sh * 0.80)),
        Emu(int(sw * 0.92)), Emu(int(sh * 0.17)),
    ).table
    headers = ["", *[f"W#{w.wafer_number}" for w in wafers]]
    rows = [
        ("N", [str(w.n) for w in wafers]),
        ("Mean", [_fmt(w.mean) for w in wafers]),
        ("Cpk", [_fmt(w.cpk, 3) for w in wafers]),
    ]
    for j, h in enumerate(headers):
        table.cell(0, j).text = h
    for i, (label, values) in enumerate(rows, start=1):
        table.cell(i, 0).text = label
        for j, v in enumerate(values, start=1):
            table.cell(i, j).text = v
    _shrink_table_font(table, 9)


def _shrink_table_font(table, size_pt: int) -> None:
    for row in table.rows:
        for cell in row.cells:
            for para in cell.text_frame.paragraphs:
                for run in para.runs:
                    run.font.size = Pt(size_pt)


# ----------------------------------------------------------------------- entrypoint
def build_report(ds: "analysis.BasicTypeDataset", output_path: Path) -> Path:
    """Build the full ``.pptx`` for a dataset and write it to ``output_path``."""
    prs = _new_presentation()
    if prs.slide_width is None:
        prs.slide_width = Emu(12192000)
        prs.slide_height = Emu(6858000)

    n_files = len(ds.wafers)
    tmpdir = Path(tempfile.mkdtemp(prefix="fe_te_"))

    try:
        # 1. Title
        _add_title_slide(prs, ds, n_files)

        # 2. Overall yield
        yields = analysis.compute_yields(ds)
        yld_png = _yield_png(yields, tmpdir / "yield.png")
        _add_yield_slide(prs, ds, yields, yld_png)

        # 3. HBIN pareto
        pareto = analysis.hbin_pareto(ds, fails_only=True)
        par_png = _pareto_png(pareto, tmpdir / "pareto.png")
        _add_chart_slide(prs, "HBIN Pareto", par_png,
                         note=None if pareto else "No fail bins found.")

        # 4. HBIN trend
        wnums, top_bins, trend = analysis.hbin_trend(ds, top_n=5)
        trd_png = _trend_png(wnums, top_bins, trend, tmpdir / "trend.png")
        _add_chart_slide(prs, "HBIN Trend", trd_png,
                         note=None if top_bins else "No fail bins found.")

        # 5..N. One slide per test number
        for number in ds.test_numbers:
            summary = analysis.summarize_test(ds, number)
            box_png = _boxplot_png(summary, tmpdir / f"box_{number}.png")
            cum_png = _cumfreq_png(summary, tmpdir / f"cum_{number}.png")
            slide = _add_test_slide(prs, summary)
            _add_test_charts(slide, prs, box_png, cum_png)
            _add_test_table(slide, prs, summary)

        output_path.parent.mkdir(parents=True, exist_ok=True)
        try:
            prs.save(str(output_path))
            return output_path
        except PermissionError:
            # Target is locked (typically open in PowerPoint). Fall back to a
            # timestamped name so the run still succeeds instead of failing outright.
            from datetime import datetime

            stamp = datetime.now().strftime("%Y%m%d_%H%M%S")
            fallback = output_path.with_name(
                f"{output_path.stem}_{stamp}{output_path.suffix}"
            )
            log.warning(
                "Output %s is locked (open in another app). Writing to %s instead.",
                output_path.name,
                fallback.name,
            )
            prs.save(str(fallback))
            return fallback
    finally:
        for png in tmpdir.glob("*.png"):
            png.unlink(missing_ok=True)
        try:
            tmpdir.rmdir()
        except OSError:
            pass


def default_output_path(ds: "analysis.BasicTypeDataset") -> Path:
    stamp = datetime.now().strftime("%Y%m%d_%H%M%S_%f")
    filename = f"{ds.p_code} \u2013 {ds.vf_code} FE Testing Results {stamp}.pptx"
    return config.OUTPUT_DIR / filename
